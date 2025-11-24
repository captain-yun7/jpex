#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_profile_slide(prs):
    """
    프로필 슬라이드 생성 (2페이지 스타일)
    왼쪽: 프로필 정보, 오른쪽: 경력 및 스킬
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 빈 슬라이드

    # 배경 워터마크 "Profile" (왼쪽 영역)
    watermark_box = slide.shapes.add_textbox(Inches(0.3), Inches(2), Inches(2), Inches(3))
    watermark_frame = watermark_box.text_frame
    watermark_frame.text = "Profile"
    watermark_para = watermark_frame.paragraphs[0]
    watermark_para.font.size = Pt(72)
    watermark_para.font.bold = True
    watermark_para.font.color.rgb = RGBColor(230, 230, 230)  # 연한 회색
    watermark_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    watermark_frame.rotation = 270  # 세로로 회전

    # 왼쪽 영역 - 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(3.5), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "미래를 설계하는 개발자"
    title_frame.word_wrap = True
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(28)
    title_para.font.color.rgb = RGBColor(100, 150, 200)

    # 부제목
    subtitle_para = title_frame.add_paragraph()
    subtitle_para.text = "이미래 입니다."
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = RGBColor(100, 150, 200)

    # 프로필 사진 플레이스홀더 (원형)
    profile_circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(1.2), Inches(2.2),
        Inches(2), Inches(2)
    )
    profile_circle.fill.solid()
    profile_circle.fill.fore_color.rgb = RGBColor(200, 220, 240)  # 연한 파란색
    profile_circle.line.color.rgb = RGBColor(100, 150, 200)
    profile_circle.line.width = Pt(3)

    # 프로필 사진 텍스트
    profile_tf = profile_circle.text_frame
    profile_tf.text = "프로필\n사진"
    profile_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    profile_p = profile_tf.paragraphs[0]
    profile_p.alignment = PP_ALIGN.CENTER
    profile_p.font.size = Pt(14)
    profile_p.font.color.rgb = RGBColor(100, 100, 100)
    profile_p.line_spacing = 1.5

    # 이름과 생년월일
    name_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.4), Inches(2.5), Inches(0.5))
    name_frame = name_box.text_frame
    name_frame.text = "이미래"
    name_para = name_frame.paragraphs[0]
    name_para.font.size = Pt(26)
    name_para.font.bold = True
    name_para.font.color.rgb = RGBColor(40, 40, 40)
    name_para.alignment = PP_ALIGN.CENTER

    birth_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.8), Inches(2.5), Inches(0.3))
    birth_frame = birth_box.text_frame
    birth_frame.text = "1995.08.28"
    birth_para = birth_frame.paragraphs[0]
    birth_para.font.size = Pt(16)
    birth_para.font.color.rgb = RGBColor(100, 100, 100)
    birth_para.alignment = PP_ALIGN.CENTER

    # 연락처 정보
    contact_y = 5.3
    contacts = [
        ("M", "010.7552.3951", RGBColor(100, 150, 200)),
        ("E", "mirae3951@gmail.com", RGBColor(100, 150, 200)),
        ("A", "서울 강서구 화곡동", RGBColor(100, 150, 200))
    ]

    for label, info, color in contacts:
        # 라벨
        label_box = slide.shapes.add_textbox(Inches(0.6), Inches(contact_y), Inches(0.3), Inches(0.3))
        label_frame = label_box.text_frame
        label_frame.text = label
        label_para = label_frame.paragraphs[0]
        label_para.font.size = Pt(14)
        label_para.font.bold = True
        label_para.font.color.rgb = color

        # 정보
        info_box = slide.shapes.add_textbox(Inches(1.0), Inches(contact_y), Inches(3), Inches(0.3))
        info_frame = info_box.text_frame
        info_frame.text = info
        info_para = info_frame.paragraphs[0]
        info_para.font.size = Pt(11)
        info_para.font.color.rgb = RGBColor(60, 60, 60)

        contact_y += 0.35

    # 오른쪽 영역 - 섹션 구분선
    divider_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(4.2), Inches(0.5),
        Inches(0.02), Inches(7)
    )
    divider_line.fill.solid()
    divider_line.fill.fore_color.rgb = RGBColor(200, 200, 200)
    divider_line.line.fill.background()

    # 오른쪽 영역 - 경력사항
    right_y = 0.5

    # 경력사항 헤더
    career_header = slide.shapes.add_textbox(Inches(4.5), Inches(right_y), Inches(2.5), Inches(0.4))
    career_header_frame = career_header.text_frame
    career_header_frame.text = "경력사항"
    career_header_para = career_header_frame.paragraphs[0]
    career_header_para.font.size = Pt(16)
    career_header_para.font.bold = True
    career_header_para.font.color.rgb = RGBColor(100, 150, 200)

    # 경력사항 내용
    career_box = slide.shapes.add_textbox(Inches(4.5), Inches(right_y + 0.35), Inches(5), Inches(0.8))
    career_frame = career_box.text_frame
    career_frame.text = "개발직 경력: 1년 11개월"
    career_p = career_frame.paragraphs[0]
    career_p.font.size = Pt(11)
    career_p.font.color.rgb = RGBColor(60, 60, 60)

    career_p1 = career_frame.add_paragraph()
    career_p1.text = "2023.05 ~ 2025.04        인피닉 ( AI팀, 연구원 )\n                                     - AI 연구소"
    career_p1.font.size = Pt(10)
    career_p1.font.color.rgb = RGBColor(80, 80, 80)
    career_p1.space_before = Pt(5)

    right_y += 1.4

    # 교육 헤더
    edu_header = slide.shapes.add_textbox(Inches(4.5), Inches(right_y), Inches(2.5), Inches(0.4))
    edu_header_frame = edu_header.text_frame
    edu_header_frame.text = "교육"
    edu_header_para = edu_header_frame.paragraphs[0]
    edu_header_para.font.size = Pt(16)
    edu_header_para.font.bold = True
    edu_header_para.font.color.rgb = RGBColor(100, 150, 200)

    # 교육 내용
    edu_box = slide.shapes.add_textbox(Inches(4.5), Inches(right_y + 0.35), Inches(5), Inches(0.6))
    edu_frame = edu_box.text_frame
    edu_frame.text = "2025.04 ~ 진행중        AWS 클라우드를 활용한\n                                     MSA 기반 자바 개발자 과정"
    edu_p = edu_frame.paragraphs[0]
    edu_p.font.size = Pt(10)
    edu_p.font.color.rgb = RGBColor(80, 80, 80)
    edu_p.line_spacing = 1.3

    right_y += 1.1

    # 자격증 헤더
    cert_header = slide.shapes.add_textbox(Inches(4.5), Inches(right_y), Inches(2.5), Inches(0.4))
    cert_header_frame = cert_header.text_frame
    cert_header_frame.text = "자격증"
    cert_header_para = cert_header_frame.paragraphs[0]
    cert_header_para.font.size = Pt(16)
    cert_header_para.font.bold = True
    cert_header_para.font.color.rgb = RGBColor(100, 150, 200)

    # 자격증 내용
    cert_box = slide.shapes.add_textbox(Inches(4.5), Inches(right_y + 0.35), Inches(5), Inches(0.8))
    cert_frame = cert_box.text_frame

    certs = [
        "- 빅데이터 분석기사",
        "- SQL (SQLD)",
        "- 데이터 분석 준전문가 (ADsP)",
        "- 컴퓨터 활용 능력 (1급)"
    ]

    for i, cert in enumerate(certs):
        if i == 0:
            cert_frame.text = cert
            cert_p = cert_frame.paragraphs[0]
        else:
            cert_p = cert_frame.add_paragraph()
            cert_p.text = cert
        cert_p.font.size = Pt(10)
        cert_p.font.color.rgb = RGBColor(80, 80, 80)
        cert_p.space_after = Pt(3)

    right_y += 1.2

    # 출간 헤더
    pub_header = slide.shapes.add_textbox(Inches(7.3), Inches(0.5), Inches(2.5), Inches(0.4))
    pub_header_frame = pub_header.text_frame
    pub_header_frame.text = "출간"
    pub_header_para = pub_header_frame.paragraphs[0]
    pub_header_para.font.size = Pt(16)
    pub_header_para.font.bold = True
    pub_header_para.font.color.rgb = RGBColor(100, 150, 200)

    # 출간 내용
    pub_box = slide.shapes.add_textbox(Inches(7.3), Inches(0.85), Inches(2.5), Inches(0.6))
    pub_frame = pub_box.text_frame
    pub_frame.text = "'자율 주행 인지기술을 위한 3차원 시맨틱 세그멘테이션'\n\n[정보통신진흥기관 조가기술동향 2115호 게재 (Page 16~29)]"
    pub_frame.word_wrap = True
    pub_p = pub_frame.paragraphs[0]
    pub_p.font.size = Pt(8)
    pub_p.font.color.rgb = RGBColor(80, 80, 80)
    pub_p.line_spacing = 1.2

    # 특허 헤더
    patent_header = slide.shapes.add_textbox(Inches(7.3), Inches(1.6), Inches(2.5), Inches(0.4))
    patent_header_frame = patent_header.text_frame
    patent_header_frame.text = "특허"
    patent_header_para = patent_header_frame.paragraphs[0]
    patent_header_para.font.size = Pt(16)
    patent_header_para.font.bold = True
    patent_header_para.font.color.rgb = RGBColor(100, 150, 200)

    # 특허 내용
    patent_box = slide.shapes.add_textbox(Inches(7.3), Inches(1.95), Inches(2.5), Inches(0.4))
    patent_frame = patent_box.text_frame
    patent_frame.text = "인공지능을 이용한 교육 모니터링 방법 및 이를\n실행하기 위하여 기록매체에 기록된 컴퓨터 프로그램"
    patent_frame.word_wrap = True
    patent_p = patent_frame.paragraphs[0]
    patent_p.font.size = Pt(8)
    patent_p.font.color.rgb = RGBColor(80, 80, 80)
    patent_p.line_spacing = 1.2

    patent_detail = patent_frame.add_paragraph()
    patent_detail.text = "\n출원 번호 : 10-2023-0155645"
    patent_detail.font.size = Pt(8)
    patent_detail.font.color.rgb = RGBColor(80, 80, 80)

    # 활동 헤더
    activity_header = slide.shapes.add_textbox(Inches(7.3), Inches(2.8), Inches(2.5), Inches(0.4))
    activity_header_frame = activity_header.text_frame
    activity_header_frame.text = "활동"
    activity_header_para = activity_header_frame.paragraphs[0]
    activity_header_para.font.size = Pt(16)
    activity_header_para.font.bold = True
    activity_header_para.font.color.rgb = RGBColor(100, 150, 200)

    # 활동 내용
    activity_box = slide.shapes.add_textbox(Inches(7.3), Inches(3.15), Inches(2.5), Inches(0.6))
    activity_frame = activity_box.text_frame
    activity_frame.text = "2023.03        산림 빅데이터 활용 해커톤 본선 진출, 동상"
    activity_p = activity_frame.paragraphs[0]
    activity_p.font.size = Pt(10)
    activity_p.font.color.rgb = RGBColor(80, 80, 80)

    activity_p2 = activity_frame.add_paragraph()
    activity_p2.text = "\n2023.02        (AI 부트캠프) 실무 프로젝트 발표회, 대상"
    activity_p2.font.size = Pt(10)
    activity_p2.font.color.rgb = RGBColor(80, 80, 80)

    # 스킬 헤더
    skill_header = slide.shapes.add_textbox(Inches(7.3), Inches(4.2), Inches(2.5), Inches(0.4))
    skill_header_frame = skill_header.text_frame
    skill_header_frame.text = "스킬"
    skill_header_para = skill_header_frame.paragraphs[0]
    skill_header_para.font.size = Pt(16)
    skill_header_para.font.bold = True
    skill_header_para.font.color.rgb = RGBColor(100, 150, 200)

    # 스킬 내용 - 카테고리별로 구분
    skill_categories = [
        ("Languages", "Python, Java, JavaScript, TypeScript", RGBColor(255, 230, 100)),
        ("Backend", "Spring Boot, Spring Cloud, Spring Security", RGBColor(150, 200, 150)),
        ("Frontend", "React, Next.js", RGBColor(150, 180, 220)),
        ("Database", "MySQL, Postgre", RGBColor(200, 150, 180)),
        ("Infra", "Docker, Kubernetes, AWS, Linux, WSL", RGBColor(220, 180, 150)),
        ("DL", "PyTorch, TensorFlow", RGBColor(180, 150, 200)),
        ("ML", "Pandas, Numpy, Selenium, Scikit-learn", RGBColor(200, 200, 150)),
        ("Tools", "GitHub, Notion, Postman, Wandb, Figma", RGBColor(180, 180, 180))
    ]

    skill_y = 4.55
    for category, skills, bg_color in skill_categories:
        # 카테고리 레이블 박스
        cat_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(7.3), Inches(skill_y),
            Inches(0.8), Inches(0.25)
        )
        cat_box.fill.solid()
        cat_box.fill.fore_color.rgb = RGBColor(240, 245, 250)
        cat_box.line.color.rgb = RGBColor(100, 150, 200)
        cat_box.line.width = Pt(1)

        cat_tf = cat_box.text_frame
        cat_tf.text = category
        cat_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        cat_p = cat_tf.paragraphs[0]
        cat_p.font.size = Pt(9)
        cat_p.font.bold = True
        cat_p.font.color.rgb = RGBColor(100, 150, 200)
        cat_p.alignment = PP_ALIGN.CENTER

        # 스킬 내용
        skill_box = slide.shapes.add_textbox(Inches(8.2), Inches(skill_y), Inches(1.7), Inches(0.25))
        skill_frame = skill_box.text_frame
        skill_frame.text = skills
        skill_frame.word_wrap = True
        skill_p = skill_frame.paragraphs[0]
        skill_p.font.size = Pt(8)
        skill_p.font.color.rgb = RGBColor(60, 60, 60)

        skill_y += 0.32

def main():
    """메인 함수"""
    # 프레젠테이션 생성
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 프로필 슬라이드 생성
    create_profile_slide(prs)

    # 저장
    output_path = '/home/k8s-admin/jpex/resumes/프로필_PPT_템플릿.pptx'
    prs.save(output_path)
    print(f"✅ 프로필 PPT 템플릿이 성공적으로 생성되었습니다: {output_path}")
    print(f"📊 총 {len(prs.slides)} 개의 슬라이드가 생성되었습니다.")
    print(f"")
    print(f"💡 슬라이드 구성:")
    print(f"   - 프로필 페이지 (2페이지 스타일)")
    print(f"")
    print(f"📝 템플릿 사용법:")
    print(f"   - 프로필 사진을 원형 플레이스홀더에 삽입하세요")
    print(f"   - 각 섹션의 내용을 본인의 정보로 수정하세요")
    print(f"   - 경력, 교육, 자격증, 스킬 등을 업데이트하세요")
    print(f"   - 레이아웃과 색상은 그대로 유지됩니다")

if __name__ == "__main__":
    main()
