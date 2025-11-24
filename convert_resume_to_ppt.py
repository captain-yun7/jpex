#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_title_slide(prs):
    """표지 슬라이드 생성"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 빈 슬라이드

    # 배경색 설정
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(41, 128, 185)  # 파란색

    # 제목
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "풀스택 개발자 & AI 엔지니어"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER

    # 부제목
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "이미래 포트폴리오"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(32)
    subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
    subtitle_para.alignment = PP_ALIGN.CENTER

def create_profile_slide(prs):
    """프로필 슬라이드 생성"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "프로필"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    # 기본 정보
    p = tf.paragraphs[0]
    p.text = "이미래 (1995년생)"
    p.font.size = Pt(20)
    p.font.bold = True
    p.space_after = Pt(10)

    # 연락처
    p = tf.add_paragraph()
    p.text = "📧 dlalfofo123@naver.com"
    p.font.size = Pt(16)
    p.space_after = Pt(5)

    p = tf.add_paragraph()
    p.text = "📱 010-7552-3951"
    p.font.size = Pt(16)
    p.space_after = Pt(15)

    # 학력
    p = tf.add_paragraph()
    p.text = "🎓 학력"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_after = Pt(5)

    p = tf.add_paragraph()
    p.text = "영남대학교 의류패션학과 (복수전공: 가족주거학과)"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(3)

    p = tf.add_paragraph()
    p.text = "2014.03 ~ 2020.02 졸업 (학점: 3.92/4.5)"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(15)

    # 경력
    p = tf.add_paragraph()
    p.text = "💼 경력"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_after = Pt(5)

    p = tf.add_paragraph()
    p.text = "인피닉 AI2 연구원 (2023.05 ~ 2025.04, 2년)"
    p.font.size = Pt(14)
    p.level = 1

def create_skills_slide(prs):
    """기술 스택 슬라이드 생성"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "기술 스택"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    # 백엔드
    p = tf.paragraphs[0]
    p.text = "Backend"
    p.font.size = Pt(20)
    p.font.bold = True
    p.space_after = Pt(8)

    p = tf.add_paragraph()
    p.text = "JAVA, Spring Boot, JPA, API, SQL, RDBMS"
    p.font.size = Pt(16)
    p.level = 1
    p.space_after = Pt(15)

    # 프론트엔드
    p = tf.add_paragraph()
    p.text = "Frontend"
    p.font.size = Pt(20)
    p.font.bold = True
    p.space_after = Pt(8)

    p = tf.add_paragraph()
    p.text = "JavaScript, TypeScript, React"
    p.font.size = Pt(16)
    p.level = 1
    p.space_after = Pt(15)

    # AI/ML
    p = tf.add_paragraph()
    p.text = "AI / Machine Learning"
    p.font.size = Pt(20)
    p.font.bold = True
    p.space_after = Pt(8)

    p = tf.add_paragraph()
    p.text = "Python, PyTorch, TensorFlow"
    p.font.size = Pt(16)
    p.level = 1
    p.space_after = Pt(15)

    # DevOps
    p = tf.add_paragraph()
    p.text = "DevOps / Cloud"
    p.font.size = Pt(20)
    p.font.bold = True
    p.space_after = Pt(8)

    p = tf.add_paragraph()
    p.text = "Docker, Kubernetes, AWS, Git"
    p.font.size = Pt(16)
    p.level = 1

def create_career_slide(prs):
    """경력 상세 슬라이드 생성"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "경력 사항"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    # 인피닉
    p = tf.paragraphs[0]
    p.text = "인피닉 - AI2 연구원 (2023.05 ~ 2025.04)"
    p.font.size = Pt(20)
    p.font.bold = True
    p.space_after = Pt(10)

    achievements = [
        "자체 모델 개발 프로젝트 진행 (Backbone, Neck, Head 연구·개발)",
        "최신 논문 연구 동향 파악 및 개발 모델 적용",
        "모델 성능을 SOTA 수준까지 향상",
        "3D 시맨틱 세그멘테이션 원고 작성",
        "KISA 지능형 CCTV 시험 인증 획득",
        "GUI 개발 및 특허 출원"
    ]

    for achievement in achievements:
        p = tf.add_paragraph()
        p.text = "• " + achievement
        p.font.size = Pt(14)
        p.level = 1
        p.space_after = Pt(5)

    # KT Group 인턴
    p = tf.add_paragraph()
    p.text = ""
    p.space_after = Pt(10)

    p = tf.add_paragraph()
    p.text = "KT Group - 뉴욕 해외 인턴 (2018.08 ~ 2019.08)"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_after = Pt(8)

    p = tf.add_paragraph()
    p.text = "• 패션 벤더 회사 Production 팀에서 1년간 인턴 근무"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(5)

    p = tf.add_paragraph()
    p.text = "• 미국 현지 트렌드 분석 및 상품 기획 지원"
    p.font.size = Pt(14)
    p.level = 1

def create_education_slide(prs):
    """교육 이력 슬라이드 생성"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "교육 이력"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    # SeSAC
    p = tf.paragraphs[0]
    p.text = "AWS 클라우드를 활용한 MSA 기반 자바 개발자 과정"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_after = Pt(5)

    p = tf.add_paragraph()
    p.text = "청년취업사관학교(SeSAC) | 2025.04 ~ 2025.11"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(10)

    topics = [
        "백엔드: Java, Spring Boot, Spring Security, Spring Cloud",
        "프론트엔드: JavaScript, TypeScript, React, Next.js",
        "인프라: Docker, Kubernetes, AWS 클라우드"
    ]

    for topic in topics:
        p = tf.add_paragraph()
        p.text = "• " + topic
        p.font.size = Pt(13)
        p.level = 2
        p.space_after = Pt(3)

    # 아시아경제
    p = tf.add_paragraph()
    p.text = ""
    p.space_after = Pt(10)

    p = tf.add_paragraph()
    p.text = "AI 및 빅데이터 분석 부트캠프"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_after = Pt(5)

    p = tf.add_paragraph()
    p.text = "아시아경제 교육원 | 2022.10 ~ 2023.02"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(10)

    topics2 = [
        "Python 데이터 전처리, 분석 및 모델 구현",
        "PyTorch, TensorFlow, Keras를 활용한 딥러닝 프로젝트",
        "NumPy, Pandas, Scikit-learn 활용 데이터 처리"
    ]

    for topic in topics2:
        p = tf.add_paragraph()
        p.text = "• " + topic
        p.font.size = Pt(13)
        p.level = 2
        p.space_after = Pt(3)

def create_certificates_slide(prs):
    """자격증 및 수상 슬라이드 생성"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "자격증 및 수상"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    # 자격증
    p = tf.paragraphs[0]
    p.text = "자격증"
    p.font.size = Pt(20)
    p.font.bold = True
    p.space_after = Pt(10)

    certificates = [
        "빅데이터분석기사 (2023.12)",
        "SQL 개발자 (2023.04)",
        "데이터분석 준전문가 (2023.03)",
        "컴퓨터활용능력 1급 (2020.05)",
        "2종보통운전면허 (2014.03)"
    ]

    for cert in certificates:
        p = tf.add_paragraph()
        p.text = "• " + cert
        p.font.size = Pt(14)
        p.level = 1
        p.space_after = Pt(5)

    # 수상
    p = tf.add_paragraph()
    p.text = ""
    p.space_after = Pt(15)

    p = tf.add_paragraph()
    p.text = "수상"
    p.font.size = Pt(20)
    p.font.bold = True
    p.space_after = Pt(10)

    p = tf.add_paragraph()
    p.text = "• 아시아경제 교육원 실무 프로젝트 발표회 대상 (2023)"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(5)

    p = tf.add_paragraph()
    p.text = "• 산림빅데이터 거래소 해커톤 동상 (2023)"
    p.font.size = Pt(14)
    p.level = 1

def create_motivation_slide(prs):
    """직무 지원동기 슬라이드 생성"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "직무 지원동기"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    p.text = "AI 연구를 넘어 서비스 역량으로"
    p.font.size = Pt(20)
    p.font.bold = True
    p.space_after = Pt(12)

    motivations = [
        "웹서비스의 영향력을 체감하며 백엔드 시스템의 중요성 인식",
        "AI 모델을 실제 사용자에게 전달하기 위한 서비스화 역량 필요성 깨달음",
        "Java Spring Boot, JPA, REST API, OAuth 등 백엔드 핵심 기술 학습",
        "Spring Cloud 기반 MSA 패턴 서비스 구현 경험",
        "Docker, Kubernetes를 통한 컨테이너 오케스트레이션 경험",
        "모니터링, CI/CD, 클라우드 환경 실습으로 확장성·안정성 역량 구축"
    ]

    for motivation in motivations:
        p = tf.add_paragraph()
        p.text = "• " + motivation
        p.font.size = Pt(13)
        p.level = 1
        p.space_after = Pt(6)

def create_experience_slide(prs):
    """준비와 경험 슬라이드 생성"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "직무를 위한 준비와 경험"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    p.text = "실무 경험과 지속적인 성장을 위한 노력"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_after = Pt(12)

    experiences = [
        "AI 연구소 2년간 근무: 시선 추정, 이상행동 탐지 등 딥러닝 모델 개발",
        "ViT, CNN 기반 모델 성능 실험 및 개선으로 문제 분석 역량 강화",
        "GUI 개발, KISA 성능 인증, 특허 출원 등 실무 산출물 창출",
        "Java Spring Boot, JPA로 데이터베이스 설계 및 REST API 구현",
        "Spring Security, OAuth 적용한 인증/인가 시스템 학습",
        "Next.js 프론트엔드 개발 및 백엔드-프론트엔드 협업 경험",
        "Docker, Kubernetes로 서비스 배포 및 관리 경험",
        "Grafana, Prometheus 활용 실시간 모니터링",
        "GitHub Actions, ArgoCD로 CI/CD 배포 자동화 구성"
    ]

    for exp in experiences:
        p = tf.add_paragraph()
        p.text = "• " + exp
        p.font.size = Pt(12)
        p.level = 1
        p.space_after = Pt(4)

def create_vision_slide(prs):
    """입사 포부 슬라이드 생성"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "입사 후 포부"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    p.text = "다음 세 가지를 약속드립니다"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_after = Pt(15)

    # 첫째
    p = tf.add_paragraph()
    p.text = "1. 안정적이고 확장 가능한 백엔드 시스템 구축"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(8)

    p = tf.add_paragraph()
    p.text = "Java Spring Framework와 JPA를 활용하여 증가하는 트래픽과 신규 기능에 유연하게 대응할 수 있는 시스템을 설계하고, 코드 품질 관리와 성능 모니터링을 체계적으로 적용"
    p.font.size = Pt(13)
    p.level = 1
    p.space_after = Pt(12)

    # 둘째
    p = tf.add_paragraph()
    p.text = "2. Frontend 팀과의 원활한 협업"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(8)

    p = tf.add_paragraph()
    p.text = "명확한 API 명세서 작성과 관리를 통해 효율적으로 소통하고, 요구사항을 정확히 파악하여 최적의 API 제공"
    p.font.size = Pt(13)
    p.level = 1
    p.space_after = Pt(12)

    # 셋째
    p = tf.add_paragraph()
    p.text = "3. 사용자에게 실질적인 가치 제공"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(8)

    p = tf.add_paragraph()
    p.text = "단순 기능 구현을 넘어 서비스 품질 향상과 사용자 만족도 증대를 목표로, 데이터 기반 개선과 성능 최적화로 서비스 경쟁력 제고"
    p.font.size = Pt(13)
    p.level = 1

def create_keywords_slide(prs):
    """키워드 3가지 슬라이드 생성"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "자신을 나타내는 키워드 3가지"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    # 가성비
    p = tf.paragraphs[0]
    p.text = "1. 가성비"
    p.font.size = Pt(20)
    p.font.bold = True
    p.space_after = Pt(8)

    p = tf.add_paragraph()
    p.text = "주어진 한정된 자원과 시간을 최대한 활용하여 성과를 극대화하며, 깊이 몰입하고 노력하여 효율적이고 탁월한 결과물을 만들어냅니다."
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(15)

    # 도전 정신
    p = tf.add_paragraph()
    p.text = "2. 도전 정신"
    p.font.size = Pt(20)
    p.font.bold = True
    p.space_after = Pt(8)

    p = tf.add_paragraph()
    p.text = "'안 하고 후회하지 말고, 하고 후회하자'라는 가치관으로 실패를 두려워하지 않고 도전하며, 과정에서 배우고 성장합니다."
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(15)

    # 꼼꼼함
    p = tf.add_paragraph()
    p.text = "3. 꼼꼼함"
    p.font.size = Pt(20)
    p.font.bold = True
    p.space_after = Pt(8)

    p = tf.add_paragraph()
    p.text = "업무, 아이디어, 일정 등을 철저히 메모하여 실수를 최소화하고, 효율적이고 체계적으로 업무를 처리합니다."
    p.font.size = Pt(14)
    p.level = 1

def create_closing_slide(prs):
    """마무리 슬라이드 생성"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경색 설정
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(52, 73, 94)

    # 감사 메시지
    thanks_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1))
    thanks_frame = thanks_box.text_frame
    thanks_frame.text = "감사합니다"
    thanks_para = thanks_frame.paragraphs[0]
    thanks_para.font.size = Pt(48)
    thanks_para.font.bold = True
    thanks_para.font.color.rgb = RGBColor(255, 255, 255)
    thanks_para.alignment = PP_ALIGN.CENTER

    # 연락처
    contact_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
    contact_frame = contact_box.text_frame
    contact_frame.text = "이미래 | dlalfofo123@naver.com | 010-7552-3951"
    contact_para = contact_frame.paragraphs[0]
    contact_para.font.size = Pt(18)
    contact_para.font.color.rgb = RGBColor(255, 255, 255)
    contact_para.alignment = PP_ALIGN.CENTER

def main():
    """메인 함수"""
    # 프레젠테이션 생성
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 슬라이드 생성
    create_title_slide(prs)
    create_profile_slide(prs)
    create_skills_slide(prs)
    create_career_slide(prs)
    create_education_slide(prs)
    create_certificates_slide(prs)
    create_motivation_slide(prs)
    create_experience_slide(prs)
    create_vision_slide(prs)
    create_keywords_slide(prs)
    create_closing_slide(prs)

    # 저장
    output_path = '/home/k8s-admin/jpex/resumes/이미래_이력서_포트폴리오.pptx'
    prs.save(output_path)
    print(f"✅ PPT 파일이 성공적으로 생성되었습니다: {output_path}")
    print(f"📊 총 {len(prs.slides)} 개의 슬라이드가 생성되었습니다.")

if __name__ == "__main__":
    main()
