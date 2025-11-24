#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_profile_slide(prs):
    """
    프로필 슬라이드 생성 (2페이지 스타일)
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경 워터마크 "Profile"
    watermark_box = slide.shapes.add_textbox(Inches(0.3), Inches(2), Inches(2), Inches(3))
    watermark_frame = watermark_box.text_frame
    watermark_frame.text = "Profile"
    watermark_para = watermark_frame.paragraphs[0]
    watermark_para.font.size = Pt(72)
    watermark_para.font.bold = True
    watermark_para.font.color.rgb = RGBColor(230, 230, 230)
    watermark_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    watermark_frame.rotation = 270

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(3.5), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "미래를 설계하는 개발자"
    title_frame.word_wrap = True
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(28)
    title_para.font.color.rgb = RGBColor(100, 150, 200)

    subtitle_para = title_frame.add_paragraph()
    subtitle_para.text = "이미래 입니다."
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = RGBColor(100, 150, 200)

    # 프로필 사진 플레이스홀더
    profile_circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(1.2), Inches(2.2),
        Inches(2), Inches(2)
    )
    profile_circle.fill.solid()
    profile_circle.fill.fore_color.rgb = RGBColor(200, 220, 240)
    profile_circle.line.color.rgb = RGBColor(100, 150, 200)
    profile_circle.line.width = Pt(3)

    profile_tf = profile_circle.text_frame
    profile_tf.text = "프로필\n사진"
    profile_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    profile_p = profile_tf.paragraphs[0]
    profile_p.alignment = PP_ALIGN.CENTER
    profile_p.font.size = Pt(14)
    profile_p.font.color.rgb = RGBColor(100, 100, 100)

    # 이름
    name_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.4), Inches(2.5), Inches(0.5))
    name_frame = name_box.text_frame
    name_frame.text = "이미래"
    name_para = name_frame.paragraphs[0]
    name_para.font.size = Pt(26)
    name_para.font.bold = True
    name_para.font.color.rgb = RGBColor(40, 40, 40)
    name_para.alignment = PP_ALIGN.CENTER

    # 생년월일
    birth_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.8), Inches(2.5), Inches(0.3))
    birth_frame = birth_box.text_frame
    birth_frame.text = "1995.08.28"
    birth_para = birth_frame.paragraphs[0]
    birth_para.font.size = Pt(16)
    birth_para.font.color.rgb = RGBColor(100, 100, 100)
    birth_para.alignment = PP_ALIGN.CENTER

    # 연락처
    contact_y = 5.3
    contacts = [
        ("M", "010.7552.3951"),
        ("E", "mirae3951@gmail.com"),
        ("A", "서울 강서구 화곡동")
    ]

    for label, info in contacts:
        label_box = slide.shapes.add_textbox(Inches(0.6), Inches(contact_y), Inches(0.3), Inches(0.3))
        label_frame = label_box.text_frame
        label_frame.text = label
        label_para = label_frame.paragraphs[0]
        label_para.font.size = Pt(14)
        label_para.font.bold = True
        label_para.font.color.rgb = RGBColor(100, 150, 200)

        info_box = slide.shapes.add_textbox(Inches(1.0), Inches(contact_y), Inches(3), Inches(0.3))
        info_frame = info_box.text_frame
        info_frame.text = info
        info_para = info_frame.paragraphs[0]
        info_para.font.size = Pt(11)
        info_para.font.color.rgb = RGBColor(60, 60, 60)

        contact_y += 0.35

    # 구분선
    divider_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(4.2), Inches(0.5),
        Inches(0.02), Inches(7)
    )
    divider_line.fill.solid()
    divider_line.fill.fore_color.rgb = RGBColor(200, 200, 200)
    divider_line.line.fill.background()

    # 오른쪽 섹션들
    right_y = 0.5

    # 경력사항
    career_header = slide.shapes.add_textbox(Inches(4.5), Inches(right_y), Inches(2.5), Inches(0.4))
    career_header_frame = career_header.text_frame
    career_header_frame.text = "경력사항"
    career_header_para = career_header_frame.paragraphs[0]
    career_header_para.font.size = Pt(16)
    career_header_para.font.bold = True
    career_header_para.font.color.rgb = RGBColor(100, 150, 200)

    career_box = slide.shapes.add_textbox(Inches(4.5), Inches(right_y + 0.35), Inches(5), Inches(0.8))
    career_frame = career_box.text_frame
    career_frame.text = "개발직 경력: 1년 11개월\n2023.05 ~ 2025.04        인피닉 ( AI팀, 연구원 )\n                                     - AI 연구소"
    career_p = career_frame.paragraphs[0]
    career_p.font.size = Pt(10)
    career_p.font.color.rgb = RGBColor(80, 80, 80)
    career_p.line_spacing = 1.3

    right_y += 1.4

    # 교육
    edu_header = slide.shapes.add_textbox(Inches(4.5), Inches(right_y), Inches(2.5), Inches(0.4))
    edu_header_frame = edu_header.text_frame
    edu_header_frame.text = "교육"
    edu_header_para = edu_header_frame.paragraphs[0]
    edu_header_para.font.size = Pt(16)
    edu_header_para.font.bold = True
    edu_header_para.font.color.rgb = RGBColor(100, 150, 200)

    edu_box = slide.shapes.add_textbox(Inches(4.5), Inches(right_y + 0.35), Inches(5), Inches(0.6))
    edu_frame = edu_box.text_frame
    edu_frame.text = "2025.04 ~ 진행중        AWS 클라우드를 활용한\n                                     MSA 기반 자바 개발자 과정"
    edu_p = edu_frame.paragraphs[0]
    edu_p.font.size = Pt(10)
    edu_p.font.color.rgb = RGBColor(80, 80, 80)

    right_y += 1.1

    # 자격증
    cert_header = slide.shapes.add_textbox(Inches(4.5), Inches(right_y), Inches(2.5), Inches(0.4))
    cert_header_frame = cert_header.text_frame
    cert_header_frame.text = "자격증"
    cert_header_para = cert_header_frame.paragraphs[0]
    cert_header_para.font.size = Pt(16)
    cert_header_para.font.bold = True
    cert_header_para.font.color.rgb = RGBColor(100, 150, 200)

    cert_box = slide.shapes.add_textbox(Inches(4.5), Inches(right_y + 0.35), Inches(5), Inches(0.8))
    cert_frame = cert_box.text_frame
    cert_frame.text = "- 빅데이터 분석기사\n- SQL (SQLD)\n- 데이터 분석 준전문가 (ADsP)\n- 컴퓨터 활용 능력 (1급)"
    cert_p = cert_frame.paragraphs[0]
    cert_p.font.size = Pt(10)
    cert_p.font.color.rgb = RGBColor(80, 80, 80)
    cert_p.line_spacing = 1.3

    # 스킬
    skill_header = slide.shapes.add_textbox(Inches(7.3), Inches(4.2), Inches(2.5), Inches(0.4))
    skill_header_frame = skill_header.text_frame
    skill_header_frame.text = "스킬"
    skill_header_para = skill_header_frame.paragraphs[0]
    skill_header_para.font.size = Pt(16)
    skill_header_para.font.bold = True
    skill_header_para.font.color.rgb = RGBColor(100, 150, 200)

    # 스킬 카테고리들
    skill_categories = [
        ("Languages", "Python, Java, JavaScript, TypeScript"),
        ("Backend", "Spring Boot, Spring Cloud, Spring Security"),
        ("Frontend", "React, Next.js"),
        ("Database", "MySQL, Postgre"),
        ("Infra", "Docker, Kubernetes, AWS, Linux, WSL"),
        ("DL", "PyTorch, TensorFlow"),
        ("ML", "Pandas, Numpy, Selenium, Scikit-learn"),
        ("Tools", "GitHub, Notion, Postman, Wandb, Figma")
    ]

    skill_y = 4.55
    for category, skills in skill_categories:
        cat_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(7.3), Inches(skill_y),
            Inches(0.8), Inches(0.25)
        )
        cat_box.fill.solid()
        cat_box.fill.fore_color.rgb = RGBColor(240, 245, 250)
        cat_box.line.color.rgb = RGBColor(100, 150, 200)

        cat_tf = cat_box.text_frame
        cat_tf.text = category
        cat_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        cat_p = cat_tf.paragraphs[0]
        cat_p.font.size = Pt(9)
        cat_p.font.bold = True
        cat_p.font.color.rgb = RGBColor(100, 150, 200)
        cat_p.alignment = PP_ALIGN.CENTER

        skill_box = slide.shapes.add_textbox(Inches(8.2), Inches(skill_y), Inches(1.7), Inches(0.25))
        skill_frame = skill_box.text_frame
        skill_frame.text = skills
        skill_frame.word_wrap = True
        skill_p = skill_frame.paragraphs[0]
        skill_p.font.size = Pt(8)
        skill_p.font.color.rgb = RGBColor(60, 60, 60)

        skill_y += 0.32

def create_project_intro_slide(prs):
    """
    프로젝트 소개 슬라이드 (4페이지 스타일)
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 왼쪽 파란색 배경
    left_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(5), Inches(7.5)
    )
    left_bg.fill.solid()
    left_bg.fill.fore_color.rgb = RGBColor(135, 180, 220)
    left_bg.line.fill.background()

    # 오른쪽 흰색 배경
    right_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(5), Inches(0),
        Inches(5), Inches(7.5)
    )
    right_bg.fill.solid()
    right_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
    right_bg.line.fill.background()

    # 제목들
    title1_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(4), Inches(0.8))
    title1_frame = title1_box.text_frame
    title1_frame.text = "프로젝트 1"
    title1_para = title1_frame.paragraphs[0]
    title1_para.font.size = Pt(36)
    title1_para.font.bold = True
    title1_para.font.color.rgb = RGBColor(255, 255, 255)

    title2_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4), Inches(1.2))
    title2_frame = title2_box.text_frame
    title2_frame.text = "유기견 입양 플랫폼"
    title2_para = title2_frame.paragraphs[0]
    title2_para.font.size = Pt(44)
    title2_para.font.bold = True
    title2_para.font.color.rgb = RGBColor(255, 255, 255)
    title2_frame.word_wrap = True

    # 프로젝트 개요
    overview_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(4), Inches(0.5))
    overview_frame = overview_box.text_frame
    overview_frame.text = "프로젝트 개요"
    overview_para = overview_frame.paragraphs[0]
    overview_para.font.size = Pt(16)
    overview_para.font.bold = True
    overview_para.font.color.rgb = RGBColor(255, 230, 130)

    period_box = slide.shapes.add_textbox(Inches(0.7), Inches(3.6), Inches(3.5), Inches(0.4))
    period_frame = period_box.text_frame
    period_frame.text = "기간 : 2025.06 ~ 2025.07"
    period_para = period_frame.paragraphs[0]
    period_para.font.size = Pt(14)
    period_para.font.color.rgb = RGBColor(255, 255, 255)

    # 목표
    goal_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(4), Inches(0.4))
    goal_frame = goal_box.text_frame
    goal_frame.text = "목표"
    goal_para = goal_frame.paragraphs[0]
    goal_para.font.size = Pt(16)
    goal_para.font.bold = True
    goal_para.font.color.rgb = RGBColor(255, 230, 130)

    goal_desc_box = slide.shapes.add_textbox(Inches(0.7), Inches(4.7), Inches(3.5), Inches(0.8))
    goal_desc_frame = goal_desc_box.text_frame
    goal_desc_frame.text = "유기동물 입양과 보호를 지원하는 기능 충실 플랫폼 개발"
    goal_desc_frame.word_wrap = True
    goal_desc_para = goal_desc_frame.paragraphs[0]
    goal_desc_para.font.size = Pt(13)
    goal_desc_para.font.color.rgb = RGBColor(255, 255, 255)

    # 사용 기술
    tech_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.7), Inches(4), Inches(0.4))
    tech_frame = tech_box.text_frame
    tech_frame.text = "사용 기술 & 모델"
    tech_para = tech_frame.paragraphs[0]
    tech_para.font.size = Pt(16)
    tech_para.font.bold = True
    tech_para.font.color.rgb = RGBColor(255, 230, 130)

    back_box = slide.shapes.add_textbox(Inches(0.7), Inches(6.1), Inches(3.5), Inches(0.5))
    back_frame = back_box.text_frame
    back_frame.text = "Back : Spring Boot, Spring Security, JPA, PostgreSQL, Docker"
    back_frame.word_wrap = True
    back_para = back_frame.paragraphs[0]
    back_para.font.size = Pt(11)
    back_para.font.color.rgb = RGBColor(255, 255, 255)

    front_box = slide.shapes.add_textbox(Inches(0.7), Inches(6.5), Inches(3.5), Inches(0.5))
    front_frame = front_box.text_frame
    front_frame.text = "Front: Next.js, Axios, Zustand, Tailwind CSS, Figma"
    front_frame.word_wrap = True
    front_para = front_frame.paragraphs[0]
    front_para.font.size = Pt(11)
    front_para.font.color.rgb = RGBColor(255, 255, 255)

    # 오른쪽 이미지 플레이스홀더
    img_placeholder = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.5), Inches(0.5),
        Inches(4), Inches(3)
    )
    img_placeholder.fill.solid()
    img_placeholder.fill.fore_color.rgb = RGBColor(255, 250, 205)
    img_placeholder.line.color.rgb = RGBColor(200, 200, 200)

    img_tf = img_placeholder.text_frame
    img_tf.text = "[프로젝트 이미지 또는 스크린샷]"
    img_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    img_p = img_tf.paragraphs[0]
    img_p.alignment = PP_ALIGN.CENTER
    img_p.font.size = Pt(14)
    img_p.font.color.rgb = RGBColor(150, 150, 150)

    # 3개의 버튼
    button_texts = ["팀 협업 체계화", "커뮤니티 기능 개발", "성능 최적화"]
    button_y_positions = [4.0, 5.0, 6.0]
    desc_texts = [
        "- 요구사항 API 명세서 수립\n- Kanban 보드 활용\n- GitHub Flow 전략",
        "- 자유게시판 개발\n- SSR + CSR 프론트 구현\n- API 테스트 진행",
        "- 초기 로딩(FCP) 지연과 SEO 문제 해결\n- 대량 엔티티 최적화"
    ]

    for text, y_pos, desc in zip(button_texts, button_y_positions, desc_texts):
        button = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(5.8), Inches(y_pos),
            Inches(1.8), Inches(0.6)
        )
        button.fill.solid()
        button.fill.fore_color.rgb = RGBColor(100, 150, 200)
        button.line.color.rgb = RGBColor(80, 130, 180)

        btn_tf = button.text_frame
        btn_tf.text = text
        btn_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        btn_p = btn_tf.paragraphs[0]
        btn_p.alignment = PP_ALIGN.CENTER
        btn_p.font.size = Pt(14)
        btn_p.font.bold = True
        btn_p.font.color.rgb = RGBColor(255, 255, 255)

        desc_box = slide.shapes.add_textbox(Inches(7.8), Inches(y_pos), Inches(2), Inches(0.6))
        desc_frame = desc_box.text_frame
        desc_frame.text = desc
        desc_frame.word_wrap = True
        desc_para = desc_frame.paragraphs[0]
        desc_para.font.size = Pt(9)
        desc_para.font.color.rgb = RGBColor(80, 80, 80)
        desc_para.line_spacing = 1.2

def create_project_detail_slide(prs):
    """
    프로젝트 상세 슬라이드 (5페이지 스타일)
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 타이틀
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "Project 01"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(20)
    title_para.font.italic = True
    title_para.font.color.rgb = RGBColor(100, 100, 100)

    # 섹션 1
    section1_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(4.5), Inches(0.5))
    section1_frame = section1_box.text_frame
    section1_frame.text = "01 협업 전략  |  주요 협업 방식"
    section1_para = section1_frame.paragraphs[0]
    section1_para.font.size = Pt(16)
    section1_para.font.bold = True
    section1_para.font.color.rgb = RGBColor(100, 150, 200)

    collab_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.6),
        Inches(4.5), Inches(3)
    )
    collab_box.fill.solid()
    collab_box.fill.fore_color.rgb = RGBColor(240, 245, 250)
    collab_box.line.color.rgb = RGBColor(150, 180, 210)

    collab_tf = collab_box.text_frame
    collab_tf.text = "[협업 프로세스 플로우 차트]\n\n01 Figma → 02 GitHub Flow → 03 Kanban → 04 Daily Scrum"
    collab_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    collab_p = collab_tf.paragraphs[0]
    collab_p.alignment = PP_ALIGN.CENTER
    collab_p.font.size = Pt(12)
    collab_p.font.color.rgb = RGBColor(100, 100, 100)

    # 섹션 2
    section2_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(4.5), Inches(0.5))
    section2_frame = section2_box.text_frame
    section2_frame.text = "02 주요 기능  |  전체 기능 및 담당 구현 기능"
    section2_para = section2_frame.paragraphs[0]
    section2_para.font.size = Pt(16)
    section2_para.font.bold = True
    section2_para.font.color.rgb = RGBColor(100, 150, 200)

    features = [
        "1. 공통 사용자 관리  - 회원가입, 인증/인가, 게정 관리",
        "2. 공통 자유 게시판  - Q&A와 커뮤니티 글·댓글 관리",
        "3. 보호소 게정  - 입양 등록 등록 관리",
        "4. 입양자 게정  - 입양 신청 관리"
    ]

    feature_y = 5.4
    for feature in features:
        feat_box = slide.shapes.add_textbox(Inches(0.7), Inches(feature_y), Inches(4), Inches(0.25))
        feat_frame = feat_box.text_frame
        feat_frame.text = feature
        feat_para = feat_frame.paragraphs[0]
        feat_para.font.size = Pt(10)
        feat_para.font.color.rgb = RGBColor(60, 60, 60)
        feature_y += 0.3

    # Backend/Frontend 박스
    tech_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.7), Inches(6.5),
        Inches(4), Inches(0.7)
    )
    tech_box.fill.solid()
    tech_box.fill.fore_color.rgb = RGBColor(255, 250, 230)
    tech_box.line.color.rgb = RGBColor(200, 180, 100)

    tech_tf = tech_box.text_frame
    tech_tf.text = "Backend: 기본적인 CRUD, REST API 설계\nFrontend: 상태관리(Zustand), Axios-REST API 통신\n페이지네이션"
    tech_p = tech_tf.paragraphs[0]
    tech_p.font.size = Pt(9)
    tech_p.font.color.rgb = RGBColor(80, 80, 80)
    tech_p.line_spacing = 1.3

    # 섹션 3 (오른쪽)
    section3_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.0), Inches(4.5), Inches(0.5))
    section3_frame = section3_box.text_frame
    section3_frame.text = "03 성능 최적화  |  이슈 해결 방식"
    section3_para = section3_frame.paragraphs[0]
    section3_para.font.size = Pt(16)
    section3_para.font.bold = True
    section3_para.font.color.rgb = RGBColor(100, 150, 200)

    screenshot_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.2), Inches(1.6),
        Inches(4.5), Inches(2)
    )
    screenshot_box.fill.solid()
    screenshot_box.fill.fore_color.rgb = RGBColor(245, 245, 245)
    screenshot_box.line.color.rgb = RGBColor(180, 180, 180)

    screenshot_tf = screenshot_box.text_frame
    screenshot_tf.text = "[개발 과정 스크린샷 이미지]\n\n예: GitHub PR, Postman 테스트 등"
    screenshot_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    screenshot_p = screenshot_tf.paragraphs[0]
    screenshot_p.alignment = PP_ALIGN.CENTER
    screenshot_p.font.size = Pt(12)
    screenshot_p.font.color.rgb = RGBColor(120, 120, 120)

    # 성능 최적화 박스들
    perf1_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.2), Inches(3.8),
        Inches(2.1), Inches(1.2)
    )
    perf1_box.fill.solid()
    perf1_box.fill.fore_color.rgb = RGBColor(110, 160, 210)
    perf1_box.line.fill.background()

    perf1_tf = perf1_box.text_frame
    perf1_tf.text = "대량 엔티티 최적화\n\n문제점\n필드 변경마다 엔티티를\n매번에서 율리고\ndirty checking 수행"
    perf1_tf.vertical_anchor = MSO_ANCHOR.TOP
    perf1_p = perf1_tf.paragraphs[0]
    perf1_p.alignment = PP_ALIGN.CENTER
    perf1_p.font.size = Pt(11)
    perf1_p.font.bold = True
    perf1_p.font.color.rgb = RGBColor(255, 255, 255)
    perf1_p.line_spacing = 1.3

    # 화살표
    arrow1 = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(7.4), Inches(4.3),
        Inches(0.5), Inches(0.3)
    )
    arrow1.fill.solid()
    arrow1.fill.fore_color.rgb = RGBColor(100, 150, 200)
    arrow1.line.fill.background()

    perf2_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(8.0), Inches(3.8),
        Inches(1.8), Inches(1.2)
    )
    perf2_box.fill.solid()
    perf2_box.fill.fore_color.rgb = RGBColor(250, 250, 250)
    perf2_box.line.color.rgb = RGBColor(100, 150, 200)
    perf2_box.line.width = Pt(2)

    perf2_tf = perf2_box.text_frame
    perf2_tf.text = "문제점\n\nCSR 방식\n초기 로딩 곤애\n→ FCP 지연,\nSEO 그레이\n인덱싱 실패"
    perf2_tf.vertical_anchor = MSO_ANCHOR.TOP
    perf2_p = perf2_tf.paragraphs[0]
    perf2_p.alignment = PP_ALIGN.CENTER
    perf2_p.font.size = Pt(10)
    perf2_p.font.bold = True
    perf2_p.font.color.rgb = RGBColor(60, 60, 60)
    perf2_p.line_spacing = 1.3

    # 개선 방법
    improve_box = slide.shapes.add_textbox(Inches(5.2), Inches(5.3), Inches(4.5), Inches(1.8))
    improve_frame = improve_box.text_frame
    improve_frame.text = "개선 방법"
    improve_p = improve_frame.paragraphs[0]
    improve_p.font.size = Pt(12)
    improve_p.font.bold = True
    improve_p.font.color.rgb = RGBColor(60, 60, 60)

    improve_desc1 = improve_frame.add_paragraph()
    improve_desc1.text = "JPQL update + @Modifying 사용 → 엔티티 로드 없이 1쿼리시 벗자 처리"
    improve_desc1.font.size = Pt(9)
    improve_desc1.font.color.rgb = RGBColor(80, 80, 80)
    improve_desc1.space_after = Pt(5)

    improve_desc2 = improve_frame.add_paragraph()
    improve_desc2.text = ">> Heap Usage 약 40% 감소 / 저리량 20% 증가"
    improve_desc2.font.size = Pt(9)
    improve_desc2.font.color.rgb = RGBColor(50, 50, 50)
    improve_desc2.font.bold = True

def main():
    """메인 함수"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 3개의 슬라이드 생성
    create_profile_slide(prs)
    create_project_intro_slide(prs)
    create_project_detail_slide(prs)

    # 저장
    output_path = '/home/k8s-admin/jpex/resumes/완전한_포트폴리오_PPT_템플릿.pptx'
    prs.save(output_path)
    print(f"✅ 완전한 포트폴리오 PPT 템플릿이 성공적으로 생성되었습니다!")
    print(f"📁 파일 경로: {output_path}")
    print(f"📊 총 {len(prs.slides)} 개의 슬라이드")
    print(f"")
    print(f"💡 슬라이드 구성:")
    print(f"   1. 프로필 페이지 (PDF 2페이지 스타일)")
    print(f"   2. 프로젝트 소개 페이지 (PDF 4페이지 스타일)")
    print(f"   3. 프로젝트 상세 페이지 (PDF 5페이지 스타일)")
    print(f"")
    print(f"📝 템플릿 활용 가이드:")
    print(f"   - 슬라이드 1: 본인의 프로필 정보로 수정")
    print(f"   - 슬라이드 2-3: 프로젝트마다 복사하여 사용")
    print(f"   - 이미지 플레이스홀더에 실제 스크린샷 삽입")
    print(f"   - 색상과 레이아웃은 그대로 유지")

if __name__ == "__main__":
    main()
